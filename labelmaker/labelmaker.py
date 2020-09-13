# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
import numpy as np
import pandas as pd # xlsxファイルの読み込み
from tqdm import tqdm # 処理経過をプログレスバーで表示

# センチメートルをPowerPoint上の距離に変換する関数
def Centis(length):
    centi = Inches(length/2.54)
    return centi

for i in range(5): # 5回1-3以外の値を入れると終了する．
    print("岩石は1，鉱物は2，化石は3を入力してください．")
    specimen_type = input()

    # 岩石標本用
    if int(specimen_type) == 1:
        # xlsxファイルを読込
        df = pd.read_excel("database.xlsx", sheet_name=0)
        # 要素を取得
        num = len(df) # 行数を取得
        no = df["no"]
        jpname = df["jpname"]
        name = df["name"]
        loc1 = df["locality1"]
        loc2 = df["locality2"]
        day = df["day"]
        month = df["month"]
        year = df["year"]
        collector = df["collector"]
        collection = df["collection"]
        
        # メッセージを表示
        print("岩石用のラベルを作成します．" + str(num) + "枚のラベルが作成されます．")

        # pptファイルを作成
        prs = Presentation()
        title_only_slide_layout = prs.slide_layouts[6] # 括弧中の番号で体裁を指定．6はblank．
        # 標本ラベルの体裁
        # 挿入する位置
        left = Centis(0.0)
        top = Centis(0.0)
        # 表の幅と高さ
        width = Centis(25.4)
        height = Centis(19.05)
        # フォントサイズ
        cell_font = 36

        # pandasで読み取った各行ごとに1ページずつラベルを作成
        for p in tqdm(range(num)):
            # 挿入する内容のリストを作成
            # # 行数と列数
            row_name = [ "No", "和名", "英名", "産地", "", "採取日", "採取者", "所蔵" ]
            contents = [ str(no[p]), jpname[p], name[p], loc1[p], loc2[p], str(year[p]) + "/" + str(month[p]) + "/" + str(day[p]), collector[p], collection[p] ]
            matrix = np.vstack([row_name, contents])
            rows = len(matrix[0])
            cols = len(matrix)
            #print(matrix[0][1])
            # スライドを挿入
            slide = prs.slides.add_slide(title_only_slide_layout)
            shapes = slide.shapes
            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Centis(5.0)
            table.columns[1].width = Centis(20.4)
            # write column headings
            for i in range(rows):
                for j in range(cols):
                    #table.cell(i, j).text = matrix[j][i]
                    cell = table.cell(i, j)
                    cell.text = matrix[j][i]
                    cell.text_frame.paragraphs[0].font.size = Pt(cell_font)
        prs.save('rock.pptx')
        break
    # 鉱物標本用
    elif int(specimen_type) == 2:
        df = pd.read_excel("database.xlsx", sheet_name=1)
        # 要素を取得
        num = len(df) # 行数を取得
        no = df["no"]
        jpname = df["jpname"]
        name = df["name"]
        formula = df["formula"]
        loc1 = df["locality1"]
        loc2 = df["locality2"]
        day = df["day"]
        month = df["month"]
        year = df["year"]
        collector = df["collector"]
        collection = df["collection"]

        # メッセージを表示
        print("鉱物用のラベルを作成します．" + str(num) + "枚のラベルが作成されます．")

        # pptファイルを作成
        prs = Presentation()
        title_only_slide_layout = prs.slide_layouts[6] # 括弧中の番号で体裁を指定．6はblank．
        # 標本ラベルの体裁
        # 挿入する位置
        left = Centis(0.0)
        top = Centis(0.0)
        # 表の幅と高さ
        width = Centis(25.4)
        height = Centis(19.05)
        # フォントサイズ
        cell_font = 36

        # pandasで読み取った各行ごとに1ページずつラベルを作成
        for p in tqdm(range(num)):
            # 挿入する内容のリストを作成
            # # 行数と列数
            row_name = [ "No", "和名", "英名", "化学式", "産地", "", "採取日", "採取者", "所蔵" ]
            contents = [ str(no[p]), jpname[p], name[p], formula[p], loc1[p], loc2[p], str(year[p]) + "/" + str(month[p]) + "/" + str(day[p]), collector[p], collection[p] ]
            matrix = np.vstack([row_name, contents])
            rows = len(matrix[0])
            cols = len(matrix)
            #print(matrix[0][1])
            # スライドを挿入
            slide = prs.slides.add_slide(title_only_slide_layout)
            shapes = slide.shapes
            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Centis(5.0)
            table.columns[1].width = Centis(20.4)
            # write column headings
            for i in range(rows):
                for j in range(cols):
                    #table.cell(i, j).text = matrix[j][i]
                    cell = table.cell(i, j)
                    cell.text = matrix[j][i]
                    cell.text_frame.paragraphs[0].font.size = Pt(cell_font)
        prs.save('mineral.pptx')
        break
    # 化石標本用
    elif int(specimen_type) == 3:
        df = pd.read_excel("database.xlsx", sheet_name=2)
        
        # 要素を取得
        num = len(df) # 行数を取得
        no = df["no"]
        jpname = df["jpname"]
        genus = df["genus"]
        species = df["species"]
        formation = df["formation"]
        loc1 = df["locality1"]
        loc2 = df["locality2"]
        day = df["day"]
        month = df["month"]
        year = df["year"]
        collector = df["collector"]
        collection = df["collection"]
        # メッセージを表示
        print("化石用のラベルを作成します．" + str(num) + "枚のラベルが作成されます．")

        # pptファイルを作成
        prs = Presentation()
        title_only_slide_layout = prs.slide_layouts[6] # 括弧中の番号で体裁を指定．6はblank．
        # 標本ラベルの体裁
        # 挿入する位置
        left = Centis(0.0)
        top = Centis(0.0)
        # 表の幅と高さ
        width = Centis(25.4)
        height = Centis(19.05)
        # フォントサイズ
        cell_font = 36

        # pandasで読み取った各行ごとに1ページずつラベルを作成
        for p in tqdm(range(num)):
            # 挿入する内容のリストを作成
            # # 行数と列数
            row_name = [ "No", "和名", "学名", "地層", "産地", "", "採取日", "採取者", "所蔵" ]
            contents = [ str(no[p]), jpname[p], genus[p] + " " + species[p], formation[p], loc1[p], loc2[p], str(year[p]) + "/" + str(month[p]) + "/" + str(day[p]), collector[p], collection[p] ]
            matrix = np.vstack([row_name, contents])
            rows = len(matrix[0])
            cols = len(matrix)
            #print(matrix[0][1])
            # スライドを挿入
            slide = prs.slides.add_slide(title_only_slide_layout)
            shapes = slide.shapes
            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Centis(5.0)
            table.columns[1].width = Centis(20.4)
            # write column headings
            for i in range(rows):
                for j in range(cols):
                    #table.cell(i, j).text = matrix[j][i]
                    cell = table.cell(i, j)
                    cell.text = matrix[j][i]
                    cell.text_frame.paragraphs[0].font.size = Pt(cell_font)
                    #cell.text_frame.paragraphs[0].font.italic = True
        prs.save('fossil.pptx')
        break
    else:
        print("正しい番号を入力してください．")
        continue